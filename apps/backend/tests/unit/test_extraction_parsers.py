"""Proves each CIS Phase 2 Prompt 3 format extractor against real
content built with the same library it parses (pypdf/python-docx/
python-pptx/openpyxl) — not hand-crafted byte fixtures, so a real
encoding bug in either direction would show up.
"""

import io

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from cerebrum.infrastructure.extraction.ocr import OCRResult
from cerebrum.infrastructure.extraction.parsers import (
    CsvExtractor,
    DocxExtractor,
    EmailExtractor,
    HtmlExtractor,
    ImageExtractor,
    PdfExtractor,
    PlainTextExtractor,
    PptxExtractor,
    XlsxExtractor,
)

pytestmark = pytest.mark.unit


def test_plain_text_extractor_returns_text_and_counts() -> None:
    result = PlainTextExtractor().extract(b"line one\nline two\n")
    assert result.text == "line one\nline two\n"
    assert result.metadata["line_count"] == 3


def test_html_extractor_strips_tags_and_finds_title() -> None:
    html = b"<html><head><title>My Page</title></head><body><p>Hello</p></body></html>"
    result = HtmlExtractor().extract(html)
    assert "Hello" in result.text
    assert result.metadata["title"] == "My Page"


def test_csv_extractor_reports_row_and_column_counts() -> None:
    csv_bytes = b"a,b,c\n1,2,3\n4,5,6\n"
    result = CsvExtractor().extract(csv_bytes)
    assert result.metadata["row_count"] == 3
    assert result.metadata["column_count"] == 3
    assert "a, b, c" in result.text


def test_pdf_extractor_recovers_page_count_and_metadata() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_blank_page(width=200, height=200)
    writer.metadata = {"/Title": "Test PDF", "/Author": "Alice"}
    buffer = io.BytesIO()
    writer.write(buffer)

    result = PdfExtractor().extract(buffer.getvalue())

    assert result.metadata["page_count"] == 2
    assert result.metadata["title"] == "Test PDF"
    assert result.metadata["author"] == "Alice"


def test_docx_extractor_recovers_paragraphs_and_properties() -> None:
    document = DocxDocument()
    document.add_paragraph("Hello world")
    document.add_paragraph("Second paragraph")
    document.core_properties.title = "My Doc"
    document.core_properties.author = "Bob"
    buffer = io.BytesIO()
    document.save(buffer)

    result = DocxExtractor().extract(buffer.getvalue())

    assert result.text == "Hello world\nSecond paragraph"
    assert result.metadata["paragraph_count"] == 2
    assert result.metadata["title"] == "My Doc"
    assert result.metadata["author"] == "Bob"


def test_pptx_extractor_recovers_slide_count_and_text() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Title Slide"
    buffer = io.BytesIO()
    presentation.save(buffer)

    result = PptxExtractor().extract(buffer.getvalue())

    assert result.metadata["slide_count"] == 1
    assert "Title Slide" in result.text


def test_xlsx_extractor_recovers_sheet_names_and_cell_values() -> None:
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Sheet1"
    sheet.append(["a", "b", "c"])
    sheet.append([1, 2, 3])
    buffer = io.BytesIO()
    workbook.save(buffer)

    result = XlsxExtractor().extract(buffer.getvalue())

    assert result.metadata["sheet_names"] == ["Sheet1"]
    assert "a, b, c" in result.text
    assert "1, 2, 3" in result.text


def test_email_extractor_recovers_subject_and_body() -> None:
    raw_email = (
        b"Subject: Quarterly Report\r\n"
        b"From: alice@example.com\r\n"
        b"To: bob@example.com\r\n"
        b"Content-Type: text/plain\r\n"
        b"\r\n"
        b"Please see the attached report.\r\n"
    )
    result = EmailExtractor().extract(raw_email)

    assert result.metadata["subject"] == "Quarterly Report"
    assert result.metadata["from"] == "alice@example.com"
    assert "Please see the attached report." in result.text


class _FakeOCREngine:
    def __init__(self, text: str) -> None:
        self._text = text

    def recognize(self, image_bytes: bytes) -> OCRResult:
        return OCRResult(text=self._text)


def test_image_extractor_uses_ocr_engine_and_reports_dimensions() -> None:
    from PIL import Image

    buffer = io.BytesIO()
    Image.new("RGB", (64, 32), color="white").save(buffer, format="PNG")

    result = ImageExtractor(_FakeOCREngine("recognized text")).extract(
        buffer.getvalue()
    )

    assert result.text == "recognized text"
    assert result.metadata == {"width": 64, "height": 32, "format": "PNG"}
