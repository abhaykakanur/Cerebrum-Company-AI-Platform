"""Multi-format text/metadata extractors — CIS Phase 2 Prompt 3's
Intelligent Document Processing Pipeline. Each extractor is a plain
synchronous callable (CPU-bound parsing; offloaded via
``asyncio.to_thread`` at the service layer —
cerebrum.application.knowledge.extraction_service — the same pattern
this codebase already uses for Argon2 in
cerebrum.application.auth.authentication_service), implementing the
:class:`TextExtractor` Protocol below.

Extraction here means recovering plain text and structural metadata
(page/slide/sheet counts, author/title, dimensions) — no normalization
beyond what each library already does (e.g. HTML tag stripping), and no
chunking/embeddings/entity extraction; see this milestone's
Non-Objectives.
"""

import csv
import io
from email import message_from_bytes
from email.message import Message
from typing import Protocol

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation
from pypdf import PdfReader

from cerebrum.infrastructure.extraction.ocr import OCREngine
from cerebrum.infrastructure.extraction.results import ExtractionResult


class TextExtractor(Protocol):
    def extract(self, content: bytes) -> ExtractionResult: ...


class PlainTextExtractor:
    """TXT and Markdown — Markdown is extracted as raw source text, not
    rendered to HTML/plain prose first; CIS Phase 2 Prompt 4's chunking
    stage (out of this milestone's scope) is the natural place for any
    markup-aware splitting.
    """

    def extract(self, content: bytes) -> ExtractionResult:
        text = content.decode("utf-8", errors="replace")
        return ExtractionResult(
            text=text,
            metadata={
                "character_count": len(text),
                "line_count": text.count("\n") + 1,
            },
        )


class HtmlExtractor:
    def extract(self, content: bytes) -> ExtractionResult:
        soup = BeautifulSoup(content, "html.parser")
        text = soup.get_text(separator="\n", strip=True)
        title = soup.title.string.strip() if soup.title and soup.title.string else None
        return ExtractionResult(text=text, metadata={"title": title})


class CsvExtractor:
    def extract(self, content: bytes) -> ExtractionResult:
        decoded = content.decode("utf-8", errors="replace")
        rows = list(csv.reader(io.StringIO(decoded)))
        text = "\n".join(", ".join(row) for row in rows)
        return ExtractionResult(
            text=text,
            metadata={
                "row_count": len(rows),
                "column_count": len(rows[0]) if rows else 0,
            },
        )


class PdfExtractor:
    def extract(self, content: bytes) -> ExtractionResult:
        reader = PdfReader(io.BytesIO(content))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        info = reader.metadata
        return ExtractionResult(
            text=text,
            metadata={
                "page_count": len(reader.pages),
                "title": info.title if info else None,
                "author": info.author if info else None,
            },
        )


class DocxExtractor:
    def extract(self, content: bytes) -> ExtractionResult:
        document = DocxDocument(io.BytesIO(content))
        paragraphs = [p.text for p in document.paragraphs]
        properties = document.core_properties
        return ExtractionResult(
            text="\n".join(paragraphs),
            metadata={
                "paragraph_count": len(paragraphs),
                "title": properties.title or None,
                "author": properties.author or None,
            },
        )


class PptxExtractor:
    def extract(self, content: bytes) -> ExtractionResult:
        presentation = Presentation(io.BytesIO(content))
        slide_texts = [
            "\n".join(shape.text for shape in slide.shapes if shape.has_text_frame)
            for slide in presentation.slides
        ]
        return ExtractionResult(
            text="\n\n".join(slide_texts),
            metadata={"slide_count": len(presentation.slides)},
        )


class XlsxExtractor:
    def extract(self, content: bytes) -> ExtractionResult:
        workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        try:
            row_texts = [
                ", ".join(str(cell) for cell in row if cell is not None)
                for sheet in workbook.worksheets
                for row in sheet.iter_rows(values_only=True)
            ]
        finally:
            workbook.close()
        return ExtractionResult(
            text="\n".join(row_texts),
            metadata={
                "sheet_count": len(workbook.sheetnames),
                "sheet_names": workbook.sheetnames,
            },
        )


class EmailExtractor:
    def extract(self, content: bytes) -> ExtractionResult:
        message: Message = message_from_bytes(content)
        body_parts = [
            self._decode_payload(part)
            for part in message.walk()
            if part.get_content_type() == "text/plain" and not part.get_filename()
        ]
        attachments = [name for part in message.walk() if (name := part.get_filename())]
        return ExtractionResult(
            text="\n".join(part for part in body_parts if part),
            metadata={
                "subject": message.get("Subject"),
                "from": message.get("From"),
                "to": message.get("To"),
                "attachment_count": len(attachments),
            },
        )

    @staticmethod
    def _decode_payload(part: Message) -> str:
        payload = part.get_payload(decode=True)
        if not isinstance(payload, bytes):
            return ""
        return payload.decode(part.get_content_charset() or "utf-8", errors="replace")


class ImageExtractor:
    """Runs the OCR abstraction
    (cerebrum.infrastructure.extraction.ocr) against the raw image
    bytes for text; Pillow is used only to validate the image and
    recover its dimensions/format, never to perform OCR itself.
    """

    def __init__(self, ocr_engine: OCREngine) -> None:
        self._ocr = ocr_engine

    def extract(self, content: bytes) -> ExtractionResult:
        with Image.open(io.BytesIO(content)) as image:
            width, height = image.size
            image_format = image.format
        result = self._ocr.recognize(content)
        return ExtractionResult(
            text=result.text,
            metadata={"width": width, "height": height, "format": image_format},
        )
